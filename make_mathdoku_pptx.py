from __future__ import annotations

import re
import sys
from dataclasses import dataclass
from pathlib import Path

import yaml
from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_AUTO_SHAPE_TYPE, MSO_CONNECTOR
from pptx.enum.text import MSO_ANCHOR, MSO_AUTO_SIZE, PP_ALIGN
from pptx.oxml.ns import qn
from pptx.oxml.xmlchemy import OxmlElement
from pptx.util import Emu, Inches, Pt

# Slide size: widescreen 16:9
SLIDE_W_IN = 13.333
SLIDE_H_IN = 7.5

# Colors
AXIS_LABEL_MAGENTA = RGBColor(200, 0, 200)
CAGE_LABEL_BLUE = RGBColor(30, 90, 200)
VALUE_GRAY = RGBColor(60, 65, 75)
DEFAULT_CANDIDATES_DARK_RED = RGBColor(139, 0, 0)

THIN_GRID_PT_DEFAULT = 1.0
THICK_BORDER_PT_DEFAULT = 6.5
FOOTER_TEXT = "@mnaoumov"

# Per-size layout profiles (4..9).
# These are intentionally *discrete* (not derived formulas) so the layout can
# be tuned to match the Mathdoku app look for each size.
#
# - grid_*_in / solve.*_in are absolute inches
# - axis_*_offset_frac are fractions of cell size
# - candidates.*_frac / cage.*_frac are fractions of cell size
LAYOUT_PROFILES: dict[int, dict] = {
    4: {
        "title_h_in": 0.85,
        "title_sz": 30,
        "meta_sz": 20,
        "grid_left_in": 0.65,
        "grid_top_in": 1.35,
        "grid_size_in": 4.75,
        "thin_pt": 1.0,
        "thick_pt": 5.0,
        "axis_font": 24,
        "axis_label_h": 0.34,
        "axis_label_w": 0.30,
        "axis_top_offset": 0.42,
        "axis_side_offset": 0.36,
        "value_font": 52,
        "candidates": {"x_frac": 0.10, "y_frac": 0.33, "w_frac": 0.80, "h_frac": 0.65, "font": 28},
        "cage": {"inset_x_frac": 0.03, "inset_y_frac": 0.02, "box_w_frac": 0.65, "box_h_frac": 0.35, "font": 28},
        "solve": {"enabled": True, "left_in": 6.20, "cols": 3, "col_w_in": 2.10, "col_gap_in": 0.25, "font": 16},
    },
    5: {
        "title_h_in": 0.70,
        "title_sz": 26,
        "meta_sz": 18,
        "grid_left_in": 0.65,
        "grid_top_in": 1.25,
        "grid_size_in": 5.20,
        "thin_pt": 1.0,
        "thick_pt": 5.0,
        "axis_font": 26,
        "axis_label_h": 0.37,
        "axis_label_w": 0.32,
        "axis_top_offset": 0.45,
        "axis_side_offset": 0.38,
        "value_font": 44,
        "candidates": {"x_frac": 0.10, "y_frac": 0.32, "w_frac": 0.80, "h_frac": 0.66, "font": 22},
        "cage": {"inset_x_frac": 0.03, "inset_y_frac": 0.02, "box_w_frac": 0.65, "box_h_frac": 0.33, "font": 24},
        "solve": {"enabled": True, "left_in": 6.55, "cols": 3, "col_w_in": 2.00, "col_gap_in": 0.25, "font": 16},
    },
    6: {
        "title_h_in": 0.65,
        "title_sz": 24,
        "meta_sz": 16,
        "grid_left_in": 0.65,
        "grid_top_in": 1.15,
        "grid_size_in": 5.70,
        "thin_pt": 1.0,
        "thick_pt": 6.5,
        "axis_font": 22,
        "axis_label_h": 0.32,
        "axis_label_w": 0.28,
        "axis_top_offset": 0.41,
        "axis_side_offset": 0.34,
        "value_font": 38,
        "candidates": {"x_frac": 0.10, "y_frac": 0.31, "w_frac": 0.80, "h_frac": 0.67, "font": 18},
        "cage": {"inset_x_frac": 0.03, "inset_y_frac": 0.02, "box_w_frac": 0.65, "box_h_frac": 0.30, "font": 22},
        "solve": {"enabled": True, "left_in": 6.85, "cols": 3, "col_w_in": 1.90, "col_gap_in": 0.25, "font": 16},
    },
    7: {
        "title_h_in": 0.55,
        "title_sz": 22,
        "meta_sz": 14,
        "grid_left_in": 0.65,
        "grid_top_in": 1.10,
        "grid_size_in": 6.05,
        "thin_pt": 1.0,
        "thick_pt": 6.5,
        "axis_font": 28,
        "axis_label_h": 0.40,
        "axis_label_w": 0.35,
        "axis_top_offset": 0.49,
        "axis_side_offset": 0.42,
        "value_font": 32,
        "candidates": {"x_frac": 0.10, "y_frac": 0.30, "w_frac": 0.80, "h_frac": 0.68, "font": 15},
        "cage": {"inset_x_frac": 0.03, "inset_y_frac": 0.02, "box_w_frac": 0.65, "box_h_frac": 0.28, "font": 20},
        "solve": {"enabled": True, "left_in": 7.05, "cols": 3, "col_w_in": 1.85, "col_gap_in": 0.25, "font": 16},
    },
    8: {
        "title_h_in": 0.55,
        "title_sz": 22,
        "meta_sz": 14,
        "grid_left_in": 0.65,
        "grid_top_in": 1.10,
        "grid_size_in": 6.20,
        "thin_pt": 1.0,
        "thick_pt": 6.5,
        "axis_font": 28,
        "axis_label_h": 0.40,
        "axis_label_w": 0.35,
        "axis_top_offset": 0.49,
        "axis_side_offset": 0.42,
        "value_font": 30,
        "candidates": {"x_frac": 0.10, "y_frac": 0.28, "w_frac": 0.80, "h_frac": 0.70, "font": 13},
        "cage": {"inset_x_frac": 0.03, "inset_y_frac": 0.02, "box_w_frac": 0.65, "box_h_frac": 0.26, "font": 18},
        "solve": {"enabled": True, "left_in": 7.15, "cols": 3, "col_w_in": 1.80, "col_gap_in": 0.25, "font": 16},
    },
    9: {
        "title_h_in": 0.55,
        "title_sz": 22,
        "meta_sz": 14,
        "grid_left_in": 0.65,
        "grid_top_in": 1.10,
        "grid_size_in": 6.30,
        "thin_pt": 1.0,
        "thick_pt": 6.5,
        "axis_font": 28,
        "axis_label_h": 0.40,
        "axis_label_w": 0.35,
        "axis_top_offset": 0.49,
        "axis_side_offset": 0.42,
        "value_font": 28,
        "candidates": {"x_frac": 0.10, "y_frac": 0.26, "w_frac": 0.80, "h_frac": 0.72, "font": 12},
        "cage": {"inset_x_frac": 0.03, "inset_y_frac": 0.02, "box_w_frac": 0.70, "box_h_frac": 0.24, "font": 16},
        "solve": {"enabled": True, "left_in": 7.25, "cols": 3, "col_w_in": 1.75, "col_gap_in": 0.25, "font": 16},
    },
}


def _profile(n: int) -> dict:
    if n not in LAYOUT_PROFILES:
        raise ValueError(f"Only sizes 4-9 are supported right now (got {n})")
    return LAYOUT_PROFILES[n]


@dataclass(frozen=True)
class Cage:
    cells: tuple[tuple[int, int], ...]  # (r, c) 0-based
    label: str


CELL_RE = re.compile(r"^([A-Z])([1-9][0-9]*)$")


def _font(run, *, name: str = "Segoe UI", size: int, bold: bool = False, rgb: RGBColor) -> None:
    run.font.name = name
    run.font.size = Pt(size)
    run.font.bold = bold
    run.font.color.rgb = rgb


def _set_slide_size(prs: Presentation) -> None:
    prs.slide_width = Inches(SLIDE_W_IN)
    prs.slide_height = Inches(SLIDE_H_IN)


def _lock_shape(shape, *, move: bool, resize: bool, rotate: bool, select: bool, text_edit: bool) -> None:
    """
    Apply DrawingML `a:spLocks` to shape.

    NOTE: `select=True` means "lock selection" (noSelect=1). Avoid for shapes VBA must access.
    """
    sp = getattr(shape, "_element", None)
    if sp is None:
        return
    cNvSpPr = sp.find(".//" + qn("p:cNvSpPr"))
    if cNvSpPr is None:
        return

    spLocks = cNvSpPr.find(qn("a:spLocks"))
    if spLocks is None:
        spLocks = OxmlElement("a:spLocks")
        cNvSpPr.append(spLocks)

    spLocks.set("noMove", "1" if move else "0")
    spLocks.set("noResize", "1" if resize else "0")
    spLocks.set("noRot", "1" if rotate else "0")
    spLocks.set("noSelect", "1" if select else "0")
    spLocks.set("noTextEdit", "1" if text_edit else "0")


def _add_title(slide, *, title: str, meta: str, n: int) -> None:
    prof = _profile(n)
    top = 0.05
    h = float(prof["title_h_in"])
    title_sz = int(prof["title_sz"])
    meta_sz = int(prof["meta_sz"])

    box = slide.shapes.add_textbox(Inches(0.2), Inches(top), Inches(SLIDE_W_IN - 0.4), Inches(h))
    tf = box.text_frame
    tf.clear()
    tf.vertical_anchor = MSO_ANCHOR.TOP

    p1 = tf.paragraphs[0]
    p1.text = title
    p1.alignment = PP_ALIGN.CENTER
    _font(p1.runs[0], size=title_sz, bold=True, rgb=RGBColor(0, 0, 0))

    p2 = tf.add_paragraph()
    p2.text = meta
    p2.alignment = PP_ALIGN.CENTER
    _font(p2.runs[0], size=meta_sz, bold=False, rgb=RGBColor(0, 0, 0))


def _add_footer(slide, *, text: str) -> None:
    box = slide.shapes.add_textbox(Inches(0.4), Inches(SLIDE_H_IN - 0.45), Inches(SLIDE_W_IN - 0.8), Inches(0.3))
    tf = box.text_frame
    tf.clear()
    p = tf.paragraphs[0]
    p.text = text
    p.alignment = PP_ALIGN.RIGHT
    _font(p.runs[0], size=14, bold=False, rgb=RGBColor(110, 120, 135))


def _add_hidden_meta(slide, *, puzzle_id: str, n: int, operations: bool) -> None:
    """
    Store metadata off-slide for VBA.
    Must NOT be locked for selection, or VBA access may error.
    """
    meta = slide.shapes.add_textbox(Inches(SLIDE_W_IN + 1.0), Inches(SLIDE_H_IN + 1.0), Inches(0.2), Inches(0.2))
    meta.name = "MATHDOKU_META"
    meta.fill.background()
    meta.line.fill.background()
    tf = meta.text_frame
    tf.clear()
    tf.word_wrap = False
    p = tf.paragraphs[0]
    p.text = (
        "mathdoku:\n"
        f"id: {puzzle_id}\n"
        f"size: {n}\n"
        f"operations: {str(bool(operations)).lower()}\n"
    )
    p.alignment = PP_ALIGN.LEFT
    if not p.runs:
        p.add_run()
    _font(p.runs[0], size=1, bold=False, rgb=RGBColor(255, 255, 255))
    _lock_shape(meta, move=True, resize=True, rotate=True, select=False, text_edit=True)


def _add_axis_labels(slide, *, grid_left: float, grid_top: float, cell_w: float, n: int) -> None:
    prof = _profile(n)
    axis_font = int(prof["axis_font"])
    label_h = float(prof["axis_label_h"])
    label_w = float(prof["axis_label_w"])
    top_offset = float(prof["axis_top_offset"])
    side_offset = float(prof["axis_side_offset"])

    top_y = grid_top - top_offset
    side_x = grid_left - side_offset

    for c in range(n):
        x = grid_left + c * cell_w
        box = slide.shapes.add_textbox(Inches(x), Inches(top_y), Inches(cell_w), Inches(label_h))
        tf = box.text_frame
        tf.clear()
        tf.margin_left = Inches(0)
        tf.margin_right = Inches(0)
        tf.margin_top = Inches(0)
        tf.margin_bottom = Inches(0)
        p = tf.paragraphs[0]
        p.text = chr(ord("A") + c)
        p.alignment = PP_ALIGN.CENTER
        _font(p.runs[0], size=axis_font, bold=True, rgb=AXIS_LABEL_MAGENTA)
        box.line.fill.background()
        box.fill.background()
        _lock_shape(box, move=True, resize=True, rotate=True, select=True, text_edit=True)

    for r in range(n):
        y = grid_top + r * cell_w + (cell_w - label_h) / 2
        box = slide.shapes.add_textbox(Inches(side_x), Inches(y), Inches(label_w), Inches(label_h))
        tf = box.text_frame
        tf.clear()
        tf.margin_left = Inches(0)
        tf.margin_right = Inches(0)
        tf.margin_top = Inches(0)
        tf.margin_bottom = Inches(0)
        tf.vertical_anchor = MSO_ANCHOR.MIDDLE
        p = tf.paragraphs[0]
        p.text = str(r + 1)
        p.alignment = PP_ALIGN.CENTER
        _font(p.runs[0], size=axis_font, bold=True, rgb=AXIS_LABEL_MAGENTA)
        box.line.fill.background()
        box.fill.background()
        _lock_shape(box, move=True, resize=True, rotate=True, select=True, text_edit=True)


def _add_cell_value_box(slide, *, left: float, top: float, size: float, name: str, font_size: int) -> None:
    box = slide.shapes.add_textbox(Inches(left), Inches(top), Inches(size), Inches(size))
    box.name = name
    box.fill.background()
    box.line.fill.background()

    tf = box.text_frame
    tf.clear()
    tf.vertical_anchor = MSO_ANCHOR.MIDDLE
    tf.auto_size = MSO_AUTO_SIZE.NONE
    tf.margin_left = Inches(0.0)
    tf.margin_right = Inches(0.0)
    tf.margin_top = Inches(0.0)
    tf.margin_bottom = Inches(0.0)
    tf.word_wrap = False

    p = tf.paragraphs[0]
    p.text = " "
    p.alignment = PP_ALIGN.CENTER
    _font(p.runs[0], size=font_size, bold=True, rgb=VALUE_GRAY)

    # Keep accessible to VBA. Prevent accidental movement.
    _lock_shape(box, move=True, resize=True, rotate=True, select=False, text_edit=False)


def _add_cell_candidates_box(
    slide,
    *,
    left: float,
    top: float,
    width: float,
    height: float,
    name: str,
    font_size: int,
    rgb: RGBColor,
) -> None:
    box = slide.shapes.add_textbox(Inches(left), Inches(top), Inches(width), Inches(height))
    box.name = name
    box.fill.background()
    box.line.fill.background()

    tf = box.text_frame
    tf.clear()
    tf.vertical_anchor = MSO_ANCHOR.BOTTOM
    tf.word_wrap = True
    tf.auto_size = MSO_AUTO_SIZE.NONE
    tf.margin_left = Inches(0.0)
    tf.margin_right = Inches(0.0)
    tf.margin_top = Inches(0.0)
    tf.margin_bottom = Inches(0.0)

    p = tf.paragraphs[0]
    p.text = " "
    p.alignment = PP_ALIGN.CENTER
    _font(p.runs[0], name="Consolas", size=font_size, bold=False, rgb=rgb)

    _lock_shape(box, move=True, resize=True, rotate=True, select=False, text_edit=False)


def _add_cell_hitbox(slide, *, left: float, top: float, size: float, name: str) -> None:
    box = slide.shapes.add_shape(MSO_AUTO_SHAPE_TYPE.RECTANGLE, Inches(left), Inches(top), Inches(size), Inches(size))
    box.name = name
    # No fill/no line: still clickable/selectable, but never obscures cell content.
    box.fill.background()
    box.line.fill.background()
    _lock_shape(box, move=True, resize=True, rotate=True, select=False, text_edit=True)


def _op_symbol(op: str) -> str:
    op = op.strip()
    return {
        "+": "+",
        "-": "−",
        "−": "−",
        "*": "x",
        "x": "x",
        "X": "x",
        "x": "x",
        "/": "/",
        "÷": "/",
    }.get(op, op)


def _parse_cell(token: str, *, n: int) -> tuple[int, int]:
    m = CELL_RE.match(token.strip().upper())
    if not m:
        raise ValueError(f"Bad cell ref: {token!r} (expected like A1)")
    col_ch, row_s = m.group(1), m.group(2)
    c = ord(col_ch) - ord("A")
    r = int(row_s) - 1
    if not (0 <= r < n and 0 <= c < n):
        raise ValueError(f"Cell out of range for {n}x{n}: {token!r}")
    return (r, c)


def _fit_font_size_for_box(*, text: str, base_pt: int, box_w_in: float, box_h_in: float, min_pt: int = 7) -> int:
    t = (text or "").strip()
    if not t:
        return base_pt
    max_w_pt = max(1.0, box_w_in * 72.0 - 2.0)
    max_h_pt = max(1.0, box_h_in * 72.0 - 1.0)
    n = max(1, len(t))
    width_based = int(max_w_pt / (0.60 * n))
    height_based = int(max_h_pt / 1.15)
    return max(min_pt, min(base_pt, width_based, height_based))


def _compute_boundaries(*, cell_to_cage: dict[tuple[int, int], int], n: int) -> tuple[list[list[bool]], list[list[bool]]]:
    v_bound: list[list[bool]] = [[False for _c in range(1, n)] for _r in range(n)]
    h_bound: list[list[bool]] = [[False for _c in range(n)] for _r in range(1, n)]
    for r in range(n):
        for c in range(1, n):
            v_bound[r][c - 1] = cell_to_cage[(r, c - 1)] != cell_to_cage[(r, c)]
    for r in range(1, n):
        for c in range(n):
            h_bound[r - 1][c] = cell_to_cage[(r - 1, c)] != cell_to_cage[(r, c)]
    return v_bound, h_bound


def _draw_segment(slide, *, x1: float, y1: float, x2: float, y2: float, width_pt: float, rgb: RGBColor) -> None:
    ln = slide.shapes.add_connector(MSO_CONNECTOR.STRAIGHT, Inches(x1), Inches(y1), Inches(x2), Inches(y2))
    ln.line.color.rgb = rgb
    ln.line.width = Pt(width_pt)
    _lock_shape(ln, move=True, resize=True, rotate=True, select=True, text_edit=True)


def _draw_thick_rect(slide, *, left: float, top: float, width: float, height: float) -> None:
    EMU_PER_INCH = 914400
    l = int(round(left * EMU_PER_INCH))
    t = int(round(top * EMU_PER_INCH))
    w = int(round(width * EMU_PER_INCH))
    h = int(round(height * EMU_PER_INCH))
    if w <= 0 or h <= 0:
        return
    rect = slide.shapes.add_shape(MSO_AUTO_SHAPE_TYPE.RECTANGLE, Emu(l), Emu(t), Emu(w), Emu(h))
    rect.fill.solid()
    rect.fill.fore_color.rgb = RGBColor(0, 0, 0)
    rect.line.fill.background()
    _lock_shape(rect, move=True, resize=True, rotate=True, select=True, text_edit=True)


def _draw_thin_internal_grid(
    slide,
    *,
    grid_left: float,
    grid_top: float,
    grid_size: float,
    n: int,
    thin_pt: float,
    v_bound: list[list[bool]],
    h_bound: list[list[bool]],
) -> None:
    cell_w = grid_size / n
    thin_rgb = RGBColor(170, 170, 170)
    for r in range(n):
        for c in range(1, n):
            if v_bound[r][c - 1]:
                continue
            x = grid_left + c * cell_w
            y1 = grid_top + r * cell_w
            y2 = y1 + cell_w
            _draw_segment(slide, x1=x, y1=y1, x2=x, y2=y2, width_pt=thin_pt, rgb=thin_rgb)
    for r in range(1, n):
        for c in range(n):
            if h_bound[r - 1][c]:
                continue
            y = grid_top + r * cell_w
            x1 = grid_left + c * cell_w
            x2 = x1 + cell_w
            _draw_segment(slide, x1=x1, y1=y, x2=x2, y2=y, width_pt=thin_pt, rgb=thin_rgb)


def _draw_cage_boundaries(
    slide,
    *,
    grid_left: float,
    grid_top: float,
    grid_size: float,
    n: int,
    cage_pt: float,
    v_bound: list[list[bool]],
    h_bound: list[list[bool]],
) -> None:
    cell_w = grid_size / n
    thick_w = cage_pt / 72.0
    inset = cage_pt / 144.0

    # vertical runs
    for c in range(1, n):
        r0 = 0
        while r0 < n:
            if not v_bound[r0][c - 1]:
                r0 += 1
                continue
            r1 = r0
            while r1 + 1 < n and v_bound[r1 + 1][c - 1]:
                r1 += 1
            x = grid_left + c * cell_w
            y1 = grid_top + r0 * cell_w
            y2 = grid_top + (r1 + 1) * cell_w
            if r0 == 0:
                y1 += inset
            if r1 == n - 1:
                y2 -= inset
            _draw_thick_rect(slide, left=x - thick_w / 2, top=y1, width=thick_w, height=(y2 - y1))
            r0 = r1 + 1

    # horizontal runs
    for r in range(1, n):
        c0 = 0
        while c0 < n:
            if not h_bound[r - 1][c0]:
                c0 += 1
                continue
            c1 = c0
            while c1 + 1 < n and h_bound[r - 1][c1 + 1]:
                c1 += 1
            y = grid_top + r * cell_w
            x1 = grid_left + c0 * cell_w
            x2 = grid_left + (c1 + 1) * cell_w
            if c0 == 0:
                x1 += inset
            if c1 == n - 1:
                x2 -= inset
            _draw_thick_rect(slide, left=x1, top=y - thick_w / 2, width=(x2 - x1), height=thick_w)
            c0 = c1 + 1


def _draw_thick_join_squares(
    slide,
    *,
    grid_left: float,
    grid_top: float,
    grid_size: float,
    n: int,
    cage_pt: float,
    v_bound: list[list[bool]],
    h_bound: list[list[bool]],
) -> None:
    cell_w = grid_size / n
    thick_w = cage_pt / 72.0

    def clamp(v: float, lo: float, hi: float) -> float:
        return lo if v < lo else hi if v > hi else v

    for vr in range(1, n):
        for vc in range(1, n):
            touch = False
            if v_bound[vr - 1][vc - 1] or v_bound[vr][vc - 1]:
                touch = True
            if h_bound[vr - 1][vc - 1] or h_bound[vr - 1][vc]:
                touch = True
            if not touch:
                continue
            x = grid_left + vc * cell_w
            y = grid_top + vr * cell_w
            left = clamp(x - thick_w / 2, grid_left, grid_left + grid_size - thick_w)
            top = clamp(y - thick_w / 2, grid_top, grid_top + grid_size - thick_w)
            _draw_thick_rect(slide, left=left, top=top, width=thick_w, height=thick_w)


def _draw_outer_border(
    slide,
    *,
    grid_left: float,
    grid_top: float,
    grid_size: float,
    cage_pt: float,
) -> None:
    thick_w = cage_pt / 72.0
    half = thick_w / 2.0
    # Top & bottom include corners
    _draw_thick_rect(slide, left=grid_left - half, top=grid_top - half, width=grid_size + thick_w, height=thick_w)
    _draw_thick_rect(
        slide,
        left=grid_left - half,
        top=grid_top + grid_size - half,
        width=grid_size + thick_w,
        height=thick_w,
    )
    # Left/right exclude corners
    _draw_thick_rect(slide, left=grid_left - half, top=grid_top + half, width=thick_w, height=max(0.0, grid_size - thick_w))
    _draw_thick_rect(
        slide,
        left=grid_left + grid_size - half,
        top=grid_top + half,
        width=thick_w,
        height=max(0.0, grid_size - thick_w),
    )


def _compute_layout(n: int, *, solve_enabled: bool, solve_cols: int, solve_min_col_w: float) -> dict:
    # Intentionally ignore solve_* args: all UI layout is hardcoded per-size.
    prof = _profile(n)
    solve = prof["solve"]
    return {
        "grid_left": float(prof["grid_left_in"]),
        "grid_top": float(prof["grid_top_in"]),
        "grid_size": float(prof["grid_size_in"]),
        "solve_enabled": bool(solve["enabled"]),
        "solve_left": float(solve["left_in"]),
    }


def load_yaml_spec(path: Path) -> dict:
    data = yaml.safe_load(path.read_text(encoding="utf-8"))
    if not isinstance(data, dict):
        raise ValueError("YAML spec must be a mapping at the top level")
    return data


def cages_from_yaml(spec: dict) -> tuple[int, list[Cage], str, str, bool, str, str | None]:
    n = int(spec.get("size", 4))
    difficulty = spec.get("difficulty", None)
    operations = bool(spec.get("operations", True))
    puzzle_id = str(spec.get("id", "")).strip()
    if not puzzle_id:
        raise ValueError("YAML must include non-empty 'id'")

    title = str(spec.get("title", "")).strip()
    if not title:
        try:
            title = f"#Mathdoku {int(puzzle_id):02d}"
        except Exception:
            title = f"#Mathdoku {puzzle_id}"

    meta = str(spec.get("meta", "")).strip()
    if not meta:
        parts = [f"Size {n}x{n}"]
        if difficulty is not None:
            parts.append(f"Difficulty {difficulty}")
        parts.append("With operators" if operations else "Without operators")
        meta = " • ".join(parts)

    out_in_yaml = spec.get("output", None)

    cages_in = spec.get("cages", None)
    if not isinstance(cages_in, list) or not cages_in:
        raise ValueError("cages must be a non-empty list")

    cages: list[Cage] = []
    for idx, item in enumerate(cages_in):
        if not isinstance(item, dict):
            raise ValueError(f"cages[{idx}] must be a mapping")
        cells_raw = item.get("cells", None)
        if not isinstance(cells_raw, list) or not cells_raw:
            raise ValueError(f"cages[{idx}].cells must be a non-empty list")
        cells = tuple(_parse_cell(str(t), n=n) for t in cells_raw)

        # `label` may be intentionally empty (e.g. singleton cages in stress tests).
        if "label" in item:
            label = str(item.get("label", "")).strip()
        else:
            label = ""
            if "value" not in item:
                raise ValueError(f"cages[{idx}] must have 'label' or 'value'")
            value = str(item["value"]).strip()
            op = item.get("op", item.get("operator", None))
            if operations and len(cells) > 1:
                if not op:
                    raise ValueError(f"cages[{idx}] multi-cell with operations=true requires op")
                label = f"{value}{_op_symbol(str(op))}"
            else:
                label = value

        cages.append(Cage(cells=cells, label=label))

    return n, cages, title, meta, operations, puzzle_id, out_in_yaml


def build_pptx(*, spec_path: Path, spec: dict) -> Path:
    (
        n,
        cages,
        title,
        meta,
        operations,
        puzzle_id,
        out_in_yaml,
    ) = cages_from_yaml(spec)

    prof = _profile(n)

    out_path = None
    if out_in_yaml:
        out_path = Path(str(out_in_yaml))
        if not out_path.is_absolute():
            out_path = (spec_path.parent / out_path).resolve()
    else:
        out_path = spec_path.with_suffix(".pptx")

    prs = Presentation()
    _set_slide_size(prs)
    slide = prs.slides.add_slide(prs.slide_layouts[6])

    layout = _compute_layout(n, solve_enabled=True, solve_cols=3, solve_min_col_w=1.0)

    _add_title(slide, title=title, meta=meta, n=n)
    _add_hidden_meta(slide, puzzle_id=puzzle_id, n=n, operations=bool(operations))

    grid_left = float(layout["grid_left"])
    grid_top = float(layout["grid_top"])
    grid_size = float(layout["grid_size"])
    cell_w = grid_size / n

    # Visuals are hardcoded in the per-size profile.
    value_font = int(prof["value_font"])
    cand_rgb: RGBColor = DEFAULT_CANDIDATES_DARK_RED
    thin_pt = float(prof["thin_pt"])
    thick_pt = float(prof["thick_pt"])

    cand_prof = prof["candidates"]
    cand_font = int(cand_prof["font"])

    # Candidates box placement (precalculated per-size).
    cand_left = float(cand_prof["x_frac"]) * cell_w
    cand_top = float(cand_prof["y_frac"]) * cell_w
    cand_w = float(cand_prof["w_frac"]) * cell_w
    cand_h = float(cand_prof["h_frac"]) * cell_w

    # Validate coverage/non-overlap
    all_cells = {(r, c) for r in range(n) for c in range(n)}
    cell_to_cage: dict[tuple[int, int], int] = {}
    for i, cage in enumerate(cages):
        for cell in cage.cells:
            if cell in cell_to_cage:
                raise ValueError(f"Cell {cell} appears in multiple cages")
            cell_to_cage[cell] = i
    missing = sorted(all_cells - set(cell_to_cage))
    if missing:
        raise ValueError(f"Spec missing {len(missing)} cell(s): {missing}")

    # Create cell content + hitboxes
    # Keep references to hitboxes so we can bring them above borders/labels later
    cell_hitboxes = []
    for r in range(n):
        for c in range(n):
            cell_left = grid_left + c * cell_w
            cell_top = grid_top + r * cell_w
            cell_ref = f"{chr(ord('A') + c)}{r+1}"
            _add_cell_value_box(slide, left=cell_left, top=cell_top, size=cell_w, name=f"VALUE_{cell_ref}", font_size=value_font)
            _add_cell_candidates_box(
                slide,
                left=cell_left + cand_left,
                top=cell_top + cand_top,
                width=cand_w,
                height=cand_h,
                name=f"CANDIDATES_{cell_ref}",
                font_size=cand_font,
                rgb=cand_rgb,
            )
            _add_cell_hitbox(slide, left=cell_left, top=cell_top, size=cell_w, name=f"CELL_{cell_ref}")
            cell_hitboxes.append((r, c, slide.shapes[-1]))

    # Borders
    v_bound, h_bound = _compute_boundaries(cell_to_cage=cell_to_cage, n=n)
    _draw_thin_internal_grid(slide, grid_left=grid_left, grid_top=grid_top, grid_size=grid_size, n=n, thin_pt=thin_pt, v_bound=v_bound, h_bound=h_bound)
    _draw_cage_boundaries(slide, grid_left=grid_left, grid_top=grid_top, grid_size=grid_size, n=n, cage_pt=thick_pt, v_bound=v_bound, h_bound=h_bound)
    _draw_thick_join_squares(slide, grid_left=grid_left, grid_top=grid_top, grid_size=grid_size, n=n, cage_pt=thick_pt, v_bound=v_bound, h_bound=h_bound)
    _draw_outer_border(slide, grid_left=grid_left, grid_top=grid_top, grid_size=grid_size, cage_pt=thick_pt)

    # Axis labels (drawn after borders so they appear on top)
    _add_axis_labels(slide, grid_left=grid_left, grid_top=grid_top, cell_w=cell_w, n=n)

    # Cage labels (one per cage, top-left cell)
    cage_prof = prof["cage"]
    inset_x = float(cage_prof["inset_x_frac"]) * cell_w
    inset_y = float(cage_prof["inset_y_frac"]) * cell_w
    cage_font = int(cage_prof["font"])
    for i, cage in enumerate(cages):
        tl = min(cage.cells, key=lambda rc: (rc[0], rc[1]))
        r, c = tl
        x = grid_left + c * cell_w + inset_x
        y = grid_top + r * cell_w + inset_y
        label_box_w = float(cage_prof["box_w_frac"]) * cell_w
        label_box_h = float(cage_prof["box_h_frac"]) * cell_w
        box = slide.shapes.add_textbox(Inches(x), Inches(y), Inches(label_box_w), Inches(label_box_h))
        box.name = f"CAGE_{i}_r{r+1}c{c+1}"
        box.fill.background()
        box.line.fill.background()
        tf = box.text_frame
        tf.clear()
        tf.vertical_anchor = MSO_ANCHOR.TOP
        tf.auto_size = MSO_AUTO_SIZE.TEXT_TO_FIT_SHAPE
        tf.margin_left = Inches(0.0)
        tf.margin_right = Inches(0.0)
        tf.margin_top = Inches(0.0)
        tf.margin_bottom = Inches(0.0)
        p = tf.paragraphs[0]
        p.text = cage.label
        p.alignment = PP_ALIGN.LEFT
        if not p.runs:
            p.add_run()
        _font(p.runs[0], size=cage_font, bold=True, rgb=CAGE_LABEL_BLUE)
        _lock_shape(box, move=True, resize=True, rotate=True, select=True, text_edit=True)

    _add_footer(slide, text=FOOTER_TEXT)

    # Solve notes columns (editable)
    if bool(layout["solve_enabled"]):
        solve_prof = prof["solve"]
        notes_left = float(solve_prof["left_in"])
        cols = int(solve_prof["cols"])
        col_w = float(solve_prof["col_w_in"])
        col_gap = float(solve_prof["col_gap_in"])
        notes_top = grid_top
        notes_h = grid_size
        for i in range(cols):
            box = slide.shapes.add_textbox(
                Inches(notes_left + i * (col_w + col_gap)),
                Inches(notes_top),
                Inches(col_w),
                Inches(notes_h),
            )
            box.name = f"SOLVE_NOTES_COL{i+1}"
            box.fill.background()
            box.line.color.rgb = RGBColor(200, 200, 200)
            box.line.width = Pt(1)
            tf = box.text_frame
            tf.clear()
            tf.vertical_anchor = MSO_ANCHOR.TOP
            tf.word_wrap = True
            tf.margin_left = Inches(0.08)
            tf.margin_right = Inches(0.08)
            tf.margin_top = Inches(0.08)
            tf.margin_bottom = Inches(0.08)
            p = tf.paragraphs[0]
            p.text = " "
            p.alignment = PP_ALIGN.LEFT
            _font(p.runs[0], size=int(solve_prof["font"]), bold=False, rgb=VALUE_GRAY)
            _lock_shape(box, move=True, resize=True, rotate=True, select=False, text_edit=False)

    # Bring hitboxes to front so clicks/Tab hit cells, but do it in a stable
    # row-major order to keep navigation predictable.
    for _r, _c, shp in sorted(cell_hitboxes, key=lambda t: (t[0], t[1])):
        try:
            shp.zorder(0)  # msoBringToFront
        except Exception:
            pass

    out_path.parent.mkdir(parents=True, exist_ok=True)
    prs.save(out_path)
    return out_path


def main(argv: list[str]) -> None:
    if len(argv) != 2:
        print("Error: this script accepts exactly ONE argument: the YAML spec path.")
        print("Usage: python make_mathdoku_pptx.py <spec.yaml>")
        raise SystemExit(2)
    if argv[1].startswith("-"):
        print("Error: CLI flags/parameters are not supported. YAML contains only the puzzle spec; UI/layout is hardcoded in the generator.")
        raise SystemExit(2)

    spec_path = Path(argv[1])
    spec = load_yaml_spec(spec_path)
    out = build_pptx(spec_path=spec_path, spec=spec)
    print(f"Wrote {out}")


if __name__ == "__main__":
    main(sys.argv)

